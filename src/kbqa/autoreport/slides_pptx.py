"""Automatic ~12-slide PPTX deck generator for the KBQA project.

Builds a presentation that maps 1:1 to ``docs/slide_deck_outline.md`` (Title,
Business Problem, Proposed RAG Solution, Architecture, Data Overview, Models &
Eval Results, Agentic Component, Deployment, Ethics/Privacy/Risks, Continual
Learning & Monitoring, Key Takeaways & Future Work, Q&A).

Design goals
------------
* **Lazy, optional heavy deps.** ``python-pptx`` and ``matplotlib`` (via the
  sibling :mod:`kbqa.autoreport.charts`) are imported *inside* functions. If
  ``python-pptx`` is missing we degrade to a Markdown deck so the autopilot
  never crashes.
* **Degrade gracefully.** Every artifact lookup is wrapped so a missing eval
  run, missing sibling module, or missing chart simply yields a slide with the
  static (design-brief) numbers and a short note — never an exception.
* **Aggregate stats only.** We surface metric summaries and the public sample
  questions; no large data dumps are written.

Public API
----------
``generate_slides(cfg, title=None, author=None, out_path=None) -> pathlib.Path``
returns the path to the written ``slides.pptx`` (or ``slides.md`` fallback).
"""

from __future__ import annotations

import json
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from ..config import AppConfig, artifacts_dir, run_dir
from ..logging_utils import get_logger, utc_stamp

logger = get_logger(__name__)

__all__ = ["generate_slides"]


# ─────────────────────────────────────────────────────────────────────────────
# Static content (authoritative numbers mirror docs/slide_deck_outline.md +
# DESIGN_BRIEF.md). Used verbatim when no live artifacts are available so the
# deck is always coherent and presentable.
# ─────────────────────────────────────────────────────────────────────────────

_THESIS = "Answer from your documents with citations and confidence — or say “I don’t know.”"
_DECK_TITLE = "Knowledge Base QA: Agentic RAG-over-Documents with Grounded, Cited Answers"

# Solution comparison (Slide 3).
_COMPARISON_HEADER = ["Criterion", "Plain LLM", "ChatKBQA", "RAG-over-docs (ours)"]
_COMPARISON_ROWS = [
    ["Grounded + citations", "No", "Yes (SPARQL)", "Yes (source spans)"],
    ["New domain cost", "Prompt only", "Per-schema fine-tune", "Drop in docs, re-embed"],
    ["Hardware", "GPU / API", "GPU + 100 GB+ RAM", "CPU-default"],
    ["Abstains", "Weak", "N/A", "Yes (“I don’t know”)"],
]

# Data stack (Slide 5).
_DATA_HEADER = ["Role", "Dataset", "Size", "License"]
_DATA_ROWS = [
    ["Reader + abstain", "rajpurkar/squad_v2", "130K / 11.9K", "CC-BY-SA-4.0"],
    ["Multi-hop eval", "hotpotqa/hotpot_qa", "90K train", "CC-BY-SA-4.0"],
    ["Retriever pairs", "sentence-transformers/natural-questions", "100K pairs", "CC-BY-SA-3.0"],
    ["Demo KB", "rag-datasets/rag-mini-wikipedia", "3,200 / 918 QA", "CC-BY-3.0"],
    ["Retrieval recall", "rag-datasets/rag-mini-bioasq", "gold passage IDs", "CC-BY-2.5"],
]

# Model stack (Slide 6).
_MODEL_HEADER = ["Stage", "Model (CPU default)", "GPU upgrade"]
_MODEL_ROWS = [
    ["Retriever", "BAAI/bge-base-en-v1.5", "(+ MiniLM fallback)"],
    ["Reranker", "ms-marco-MiniLM-L-6-v2", "bge-reranker-v2-m3"],
    ["Extractive reader", "deepset/roberta-base-squad2", "deberta-v3-large-squad2"],
    ["Generative reader", "google/flan-t5-base", "google/flan-t5-large"],
]

# Risk -> Mitigation (Slide 9).
_RISK_HEADER = ["Risk", "Mitigation"]
_RISK_ROWS = [
    ["Hallucination", "Grounded reader + faithfulness gate + extractive null-score → abstain"],
    ["Prompt-injection", "Treat context as data, citation-only output, reject unsupported claims"],
    ["Privacy", "Download-on-demand; no large data committed"],
    ["Licensing", "Per-source license tracking; trivia_qa / MS MARCO flagged for legal"],
    ["Cross-version index", "manifest model_version assert; blue/green swap"],
]


def _fmt(v: Any) -> str:
    """Render a metric value compactly; tolerate ``None``/non-numeric."""
    if v is None:
        return "n/a"
    try:
        return "{:.3f}".format(float(v)).rstrip("0").rstrip(".")
    except (TypeError, ValueError):
        return str(v)


def _slide_specs(art: Dict[str, Any], title: str, author: str) -> List[Dict[str, Any]]:
    """Return the ordered list of slide specifications.

    Each spec is a dict: ``{kind, title, bullets, table?, chart_key?, note?}``.
    Pulls live metrics from ``art`` (the loaded-artifacts dict) when present,
    otherwise falls back to the static design-brief numbers.
    """
    metrics = _metric_summary(art)
    versions = _model_versions(art)
    examples = _sample_questions(art)

    # Format a metric line that prefers live numbers but is honest when absent.
    rec = metrics.get("recall_hybrid")
    f1 = metrics.get("answer_f1")
    abst = metrics.get("abstain_recall")
    metric_bullets = []
    if metrics.get("_live"):
        metric_bullets.append(
            "Live eval (n={n}, k={k}): Recall@{k}={rec}, Answer F1={f1}, Abstain-recall={abst}".format(
                n=metrics.get("n_questions", "?"), k=metrics.get("k", 5),
                rec=_fmt(rec), f1=_fmt(f1), abst=_fmt(abst))
        )
    else:
        metric_bullets.append(
            "Metrics tracked: Recall@k / NDCG@10 / MRR@10; reader EM / F1 + NoAns-F1; faithfulness; citation accuracy; abstain-rate; latency"
        )
        metric_bullets.append("No eval.json found under runs/ — run `kbqa eval` to populate live numbers")
    metric_bullets.append(
        "Baseline-to-beat: BM25 + zero-shot reader (no rerank, no agent); full stack must improve EM/F1, faithfulness, citation accuracy"
    )

    specs: List[Dict[str, Any]] = [
        # 1 — Title
        {
            "kind": "title",
            "title": _DECK_TITLE,
            "subtitle": "{} · {}\nThesis: {}\nCPU-default · zero paid API · open-source models only".format(
                author, title, _THESIS),
        },
        # 2 — Business Problem
        {
            "kind": "bullets",
            "title": "Business Problem & Motivation",
            "bullets": [
                "Enterprise knowledge lives in PDFs, wikis, tickets — not curated databases; staff waste time hunting answers",
                "Plain LLMs hallucinate and give no provenance — unacceptable for compliance, support, internal knowledge",
                "Need: trustworthy answers that cite the source span and abstain when the corpus lacks support, deployable cheaply",
                "Success: high answer F1 and faithfulness, auditable citations, safe abstention, sub-second CPU latency",
            ],
        },
        # 3 — Proposed RAG Solution (table)
        {
            "kind": "table",
            "title": "Proposed RAG Solution (vs Plain LLM / ChatKBQA)",
            "bullets": [
                "Plain LLM: fluent but ungrounded, stale, no citations → hallucination risk",
                "ChatKBQA: accurate & auditable but needs ~50 GB Freebase, GPU, per-schema fine-tune — impractical for general docs",
                "Ours — RAG-over-documents: drop in docs → re-embed (no labels), CPU-default, citations, safe abstention",
            ],
            "table": {"header": _COMPARISON_HEADER, "rows": _COMPARISON_ROWS},
        },
        # 4 — Architecture (diagram chart if available)
        {
            "kind": "chart",
            "title": "System Architecture",
            "bullets": [
                "Pipeline: ingest/chunk/index → analyze → retrieve → rerank → sufficiency → grounded generate → faithfulness → answer/abstain",
                "Hybrid retrieval: BM25 (rank_bm25) ⊕ BGE dense, fused via RRF; cross-encoder rerank top-50 → top-5",
                "Two decision gates: sufficiency (CRAG loop) and faithfulness (Self-RAG gate → abstain)",
                "Every stage runs on CPU by default; GPU is a speed/accuracy knob, not a requirement",
            ],
            "chart_keys": ["architecture", "pipeline", "architecture_diagram"],
        },
        # 5 — Data Overview (table)
        {
            "kind": "table",
            "title": "Data Overview",
            "bullets": [
                "Reader: squad_v2 (~50K unanswerable questions power abstention); hotpot_qa (distractor) for multi-hop eval",
                "Retriever pairs: natural-questions for the bi-encoder fine-tune",
                "Demo/eval KB: rag-mini-wikipedia; rag-mini-bioasq gives gold passage IDs for retrieval-recall",
                "Governance: IDs verified live on HF Hub; licenses tracked; download-on-demand (no large data committed)",
            ],
            "table": {"header": _DATA_HEADER, "rows": _DATA_ROWS},
        },
        # 6 — Models & Eval Results (metrics table + chart)
        {
            "kind": "chart_table",
            "title": "Models & Evaluation Results",
            "bullets": metric_bullets,
            "table": {"header": _MODEL_HEADER, "rows": _MODEL_ROWS},
            "chart_keys": ["eval", "metrics", "eval_bars", "baseline_vs_full", "results"],
        },
        # 7 — Agentic component
        {
            "kind": "chart",
            "title": "Agentic AI Component (CRAG / Self-RAG)",
            "bullets": [
                "Deterministic state machine: query rewrite/decompose + Corrective-RAG loop + Self-RAG reflection; heuristic default + optional local LLM brain",
                "3 decision points: route simple/multi-hop/unanswerable; sufficiency loop (max_iterations={mi}, TAU_HIGH={th}, TAU_LOW={tl}); faithfulness gate → abstain".format(
                    mi=examples.get("max_iterations", 3), th=examples.get("tau_high", 0.55), tl=examples.get("tau_low", 0.15)),
                "Worked example: “Which university did the founder of SpaceX attend, and what year was it established?” → decompose into 3 dependent sub-questions; SQ2 AMBIGUOUS (0.34) → expand → SUFFICIENT (0.63) → cited answer (conf 0.90)",
                "Safety: if the founding year is missing, SQ3 stays insufficient after max_iterations → abstains rather than fabricate",
            ],
            "chart_keys": ["agent", "state_machine", "agent_trace", "crag"],
        },
        # 8 — Deployment
        {
            "kind": "bullets",
            "title": "Deployment Overview",
            "bullets": [
                "FastAPI service: /health, /ingest, /search, /ask, /batch, /metrics; plus a Gradio demo UI calling /ask",
                "FAISS persistence: kb.index + meta.parquet + manifest.json; load asserts manifest.model_version == MODEL_VERSION; blue/green swap",
                "Packaging: Docker / HF Space on port 7860; model_versions echoed in every response and /metrics",
                "Latency: /ask extractive ~350–800 ms p50/p95 on CPU; ONNX int8 (2–4×), HNSW ANN, rerank only top-50, LRU query cache",
            ],
        },
        # 9 — Ethics / Privacy / Risks (table)
        {
            "kind": "table",
            "title": "Ethics, Privacy & Risks",
            "bullets": [
                "Hallucination: grounded reader + faithfulness entailment gate + extractive null-score → abstain instead of fabricate",
                "Prompt-injection: treat retrieved context as data not commands; citation-only output; faithfulness gate rejects unsupported claims",
                "Privacy & licensing: download-on-demand; per-source license tracking; trivia_qa / MS MARCO flagged for legal review",
            ],
            "table": {"header": _RISK_HEADER, "rows": _RISK_ROWS},
        },
        # 10 — Continual Learning & Monitoring
        {
            "kind": "bullets",
            "title": "Continual Learning & Monitoring",
            "bullets": [
                "Continual ingestion: append-only add_with_ids + SHA-256 dedup (idempotent re-ingest); deletes via tombstone + periodic rebuild",
                "Retriever loop: mine hard negatives → train → rebuild index → re-mine with fine-tuned retriever → retrain (2 rounds)",
                "Monitoring: /metrics (Prometheus) exposes p50/p95 latency, cache-hit, abstain-rate, index size, request counts; watch for drift",
                "Versioning: MODEL_VERSION pins encoder + reranker + reader + index together; blue/green swap on any model change",
            ],
        },
        # 11 — Key Takeaways & Future Work
        {
            "kind": "bullets",
            "title": "Key Takeaways & Future Work",
            "bullets": [
                "Agentic RAG delivers grounded, cited, abstaining answers on CPU with zero paid API; the two gates make it production-safe",
                "Hybrid + rerank + agent-loop stack measurably beats the BM25 baseline on retrieval, F1, faithfulness, citation accuracy",
                "Future: wire the verified NLI faithfulness model; enable GPU upgrades (bge-reranker-v2-m3, deberta-v3-large, flan-t5-large)",
                "Future: add a kg_query (ChatKBQA-style) backend; longer-context encoders; optional local instruct-LLM brain (Qwen2.5-1.5B)",
            ],
        },
        # 12 — Q&A
        {
            "kind": "closing",
            "title": "Q&A",
            "bullets": [
                "Recap: “Grounded answers with citations and confidence — or an honest ‘I don’t know.’”",
                "Pre-empt: Why not a bigger LLM? How is abstention enforced? CPU latency numbers? License posture?",
                "Pointers: live Gradio demo, /ask trace output (per-decision audit log), DESIGN_BRIEF.md for full IDs and metrics",
                "Model versions: {v}".format(v=", ".join(versions) if versions else "v1"),
            ],
            "examples": examples.get("questions", []),
        },
    ]
    return specs


# ─────────────────────────────────────────────────────────────────────────────
# Artifact loading (defensive)
# ─────────────────────────────────────────────────────────────────────────────

def _try_load_all_artifacts(cfg: AppConfig) -> Dict[str, Any]:
    """Use the sibling ``load_all_artifacts`` if it exists; never raise."""
    try:  # sibling module may be written by a parallel task; import lazily.
        from .artifacts import load_all_artifacts  # type: ignore
    except Exception:
        try:
            from . import load_all_artifacts  # type: ignore  # re-exported variant
        except Exception:
            load_all_artifacts = None  # type: ignore
    if load_all_artifacts is not None:  # type: ignore
        try:
            art = load_all_artifacts(cfg)  # type: ignore
            if isinstance(art, dict):
                logger.info("Loaded artifacts via load_all_artifacts() (%d keys)", len(art))
                return art
        except Exception as exc:  # pragma: no cover - defensive
            logger.warning("load_all_artifacts() failed (%s); falling back to direct read", exc)
    return _fallback_load_artifacts(cfg)


def _fallback_load_artifacts(cfg: AppConfig) -> Dict[str, Any]:
    """Read the latest ``eval-*/eval.json`` directly from runs/ as a fallback."""
    art: Dict[str, Any] = {}
    try:
        rd = run_dir()
        eval_dirs = sorted(
            (p for p in rd.glob("eval-*") if (p / "eval.json").exists()),
            key=lambda p: p.name,
        )
        if eval_dirs:
            latest = eval_dirs[-1] / "eval.json"
            art["eval"] = json.loads(latest.read_text(encoding="utf-8"))
            logger.info("Loaded latest eval artifact: %s", latest)
        else:
            logger.info("No eval-*/eval.json under %s; using static deck numbers", rd)
    except Exception as exc:  # pragma: no cover - defensive
        logger.warning("Direct artifact read failed (%s); using static numbers", exc)
    return art


def _metric_summary(art: Dict[str, Any]) -> Dict[str, Any]:
    """Normalise whatever eval shape is present into a flat metric dict."""
    out: Dict[str, Any] = {"_live": False}
    ev = art.get("eval") if isinstance(art, dict) else None
    if not isinstance(ev, dict):
        return out
    try:
        k = ev.get("k", 5)
        summary = ev.get("summary", {}) if isinstance(ev.get("summary"), dict) else {}
        retr = ev.get("retrieval", {}) if isinstance(ev.get("retrieval"), dict) else {}
        ans = ev.get("answer", {}) if isinstance(ev.get("answer"), dict) else {}
        abst = ev.get("abstention", {}) if isinstance(ev.get("abstention"), dict) else {}
        rec_key = "recall@{}_hybrid".format(k)
        out.update({
            "_live": True,
            "k": k,
            "n_questions": ev.get("n_questions"),
            "recall_hybrid": summary.get(rec_key, retr.get(rec_key)),
            "answer_f1": summary.get("answer_f1", ans.get("f1")),
            "abstain_recall": summary.get("abstain_recall", abst.get("abstain_recall")),
        })
    except Exception as exc:  # pragma: no cover - defensive
        logger.warning("Could not parse eval summary (%s)", exc)
        out["_live"] = False
    return out


def _model_versions(art: Dict[str, Any]) -> List[str]:
    """Best-effort list of model-version strings from registry/artifacts."""
    versions: List[str] = []
    try:
        from ..models.model_registry import load_model_metadata, resolve_latest  # lazy
        from ..config import model_dir
        for sub in ("retriever", "reader", "generator", "reranker"):
            try:
                base = model_dir() / sub
                if not base.exists():
                    continue
                latest = resolve_latest(base)
                meta = load_model_metadata(latest)
                v = meta.get("model_version") or meta.get("version")
                if v:
                    versions.append("{}:{}".format(sub, v))
            except Exception:
                continue
    except Exception:
        pass
    return versions


def _sample_questions(art: Dict[str, Any]) -> Dict[str, Any]:
    """Public sample questions + agent thresholds for the closing slide."""
    info: Dict[str, Any] = {"questions": []}
    try:
        from ..data.samples import SAMPLE_QA
        for item in SAMPLE_QA[:4]:
            q = item.get("question")
            if q:
                tag = "abstain" if item.get("answer") is None else "answer"
                info["questions"].append("[{}] {}".format(tag, q))
    except Exception:
        pass
    return info


# ─────────────────────────────────────────────────────────────────────────────
# Charts (defensive)
# ─────────────────────────────────────────────────────────────────────────────

def _build_charts(cfg: AppConfig, art: Dict[str, Any], out_dir: Path) -> Dict[str, Path]:
    """Render charts via the sibling ``autoreport.charts`` module if present.

    Returns a mapping ``{chart_name: png_path}``. Always returns a dict; on any
    failure (missing module, missing matplotlib) returns what was produced so
    far (possibly empty), and the deck simply omits the missing visuals.
    """
    charts: Dict[str, Path] = {}
    try:
        from . import charts as charts_mod  # type: ignore
    except Exception as exc:
        logger.info("autoreport.charts unavailable (%s); slides will omit charts", exc)
        return charts

    out_dir.mkdir(parents=True, exist_ok=True)

    # Try a one-shot "make everything" entrypoint first (most likely API).
    for fn_name in ("generate_charts", "make_charts", "build_charts", "render_all", "all_charts"):
        fn = getattr(charts_mod, fn_name, None)
        if callable(fn):
            try:
                produced = _call_chart_fn(fn, cfg, art, out_dir)
                charts.update(_coerce_chart_map(produced))
                if charts:
                    logger.info("Charts via charts.%s() -> %d image(s)", fn_name, len(charts))
                    return charts
            except Exception as exc:  # pragma: no cover - defensive
                logger.warning("charts.%s() failed (%s)", fn_name, exc)

    # Otherwise probe individual chart factory functions by convention.
    for fn_name in dir(charts_mod):
        if not (fn_name.startswith("chart_") or fn_name.startswith("plot_") or fn_name.startswith("make_")):
            continue
        fn = getattr(charts_mod, fn_name, None)
        if not callable(fn):
            continue
        try:
            produced = _call_chart_fn(fn, cfg, art, out_dir)
            for name, path in _coerce_chart_map(produced).items():
                charts[name or fn_name] = path
        except Exception:
            continue
    if charts:
        logger.info("Charts via individual factories -> %d image(s)", len(charts))
    return charts


def _call_chart_fn(fn, cfg, art, out_dir):
    """Call a chart function tolerating a few common signatures."""
    import inspect
    try:
        sig = inspect.signature(fn)
        nparams = len([p for p in sig.parameters.values()
                       if p.kind in (p.POSITIONAL_ONLY, p.POSITIONAL_OR_KEYWORD)])
    except (ValueError, TypeError):
        nparams = 1
    attempts: List[Tuple] = []
    if nparams >= 3:
        attempts.append((cfg, art, out_dir))
    if nparams >= 2:
        attempts.append((cfg, out_dir))
        attempts.append((art, out_dir))
    attempts.append((cfg,))
    attempts.append(())
    last_exc: Optional[Exception] = None
    for args in attempts:
        try:
            return fn(*args)
        except TypeError as exc:
            last_exc = exc
            continue
    if last_exc:
        raise last_exc
    return None


def _coerce_chart_map(produced) -> Dict[str, Path]:
    """Normalise a chart function's return into ``{name: Path}``."""
    out: Dict[str, Path] = {}
    if produced is None:
        return out
    if isinstance(produced, dict):
        for k, v in produced.items():
            p = _as_existing_png(v)
            if p is not None:
                out[str(k)] = p
        return out
    if isinstance(produced, (list, tuple, set)):
        for v in produced:
            p = _as_existing_png(v)
            if p is not None:
                out[p.stem] = p
        return out
    p = _as_existing_png(produced)
    if p is not None:
        out[p.stem] = p
    return out


def _as_existing_png(v) -> Optional[Path]:
    try:
        p = Path(v)
        if p.exists() and p.suffix.lower() in (".png", ".jpg", ".jpeg", ".svg"):
            return p
    except (TypeError, ValueError):
        return None
    return None


def _match_chart(charts: Dict[str, Path], keys: List[str]) -> Optional[Path]:
    """Find the first chart whose name contains any of the requested keys."""
    if not charts:
        return None
    lowered = {name.lower(): path for name, path in charts.items()}
    for key in keys:
        k = key.lower()
        for name, path in lowered.items():
            if k in name:
                return path
    return None


# ─────────────────────────────────────────────────────────────────────────────
# PPTX rendering
# ─────────────────────────────────────────────────────────────────────────────

def _render_pptx(specs: List[Dict[str, Any]], charts: Dict[str, Path], out_path: Path) -> Path:
    """Render the slide specs to a .pptx using python-pptx (already imported)."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    accent = RGBColor(0x1F, 0x4E, 0x79)
    muted = RGBColor(0x55, 0x55, 0x55)
    light = RGBColor(0xF2, 0xF6, 0xFB)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    SW, SH = prs.slide_width, prs.slide_height

    def _add_title(slide, text: str, top=Inches(0.35), size=28):
        box = slide.shapes.add_textbox(Inches(0.6), top, SW - Inches(1.2), Inches(1.0))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = True
        run.font.color.rgb = accent
        return box

    def _add_bullets(slide, bullets: List[str], top, height, left=Inches(0.6),
                     width=None, size=15):
        width = width or (SW - Inches(1.2))
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        first = True
        for b in bullets:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            run = p.add_run()
            run.text = "• " + b
            run.font.size = Pt(size)
            run.font.color.rgb = muted
            p.space_after = Pt(6)
        return box

    def _add_table(slide, header: List[str], rows: List[List[str]], top, left=Inches(0.6),
                   width=None, height=None):
        width = width or (SW - Inches(1.2))
        nrows = len(rows) + 1
        ncols = len(header)
        height = height or Inches(0.4 * nrows)
        gf = slide.shapes.add_table(nrows, ncols, left, top, width, height)
        table = gf.table
        for c, htext in enumerate(header):
            cell = table.cell(0, c)
            cell.text = str(htext)
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.LEFT
                for run in para.runs:
                    run.font.size = Pt(12)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            cell.fill.solid()
            cell.fill.fore_color.rgb = accent
        for r, row in enumerate(rows, start=1):
            for c in range(ncols):
                cell = table.cell(r, c)
                cell.text = str(row[c]) if c < len(row) else ""
                cell.fill.solid()
                cell.fill.fore_color.rgb = light if r % 2 else RGBColor(0xFF, 0xFF, 0xFF)
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(11)
                        run.font.color.rgb = muted
        return gf

    def _add_chart_image(slide, path: Path, left, top, width):
        try:
            slide.shapes.add_picture(str(path), left, top, width=width)
            return True
        except Exception as exc:  # pragma: no cover - defensive
            logger.warning("Could not embed chart %s (%s)", path, exc)
            return False

    def _add_footer(slide, idx, total):
        box = slide.shapes.add_textbox(SW - Inches(2.2), SH - Inches(0.5), Inches(2.0), Inches(0.3))
        p = box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = "kbqa · {}/{}".format(idx, total)
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)

    total = len(specs)
    for idx, spec in enumerate(specs, start=1):
        slide = prs.slides.add_slide(blank)
        kind = spec.get("kind", "bullets")

        if kind == "title":
            band = slide.shapes.add_textbox(Inches(0.6), Inches(2.0), SW - Inches(1.2), Inches(1.6))
            tf = band.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = spec["title"]
            run.font.size = Pt(34)
            run.font.bold = True
            run.font.color.rgb = accent
            sub = slide.shapes.add_textbox(Inches(0.6), Inches(4.0), SW - Inches(1.2), Inches(2.2))
            stf = sub.text_frame
            stf.word_wrap = True
            for i, line in enumerate(spec.get("subtitle", "").split("\n")):
                pp = stf.paragraphs[0] if i == 0 else stf.add_paragraph()
                r = pp.add_run()
                r.text = line
                r.font.size = Pt(16)
                r.font.color.rgb = muted
            chip = slide.shapes.add_textbox(Inches(0.6), Inches(6.3), SW - Inches(1.2), Inches(0.6))
            cp = chip.text_frame.paragraphs[0]
            cr = cp.add_run()
            cr.text = "Grounded · Cited · Abstains"
            cr.font.size = Pt(14)
            cr.font.bold = True
            cr.font.color.rgb = accent
            _add_footer(slide, idx, total)
            continue

        _add_title(slide, spec["title"])
        bullets = spec.get("bullets", [])
        chart_path = _match_chart(charts, spec.get("chart_keys", [])) if charts else None

        if kind == "table" and spec.get("table"):
            if bullets:
                _add_bullets(slide, bullets, top=Inches(1.3), height=Inches(1.8), size=14)
                tbl_top = Inches(3.3)
            else:
                tbl_top = Inches(1.6)
            _add_table(slide, spec["table"]["header"], spec["table"]["rows"], top=tbl_top)

        elif kind == "chart_table":
            # Left: bullets + chart; right: model table.
            half = (SW - Inches(1.5)) / 2
            _add_bullets(slide, bullets, top=Inches(1.3), height=Inches(2.2),
                         left=Inches(0.6), width=half, size=13)
            if chart_path:
                _add_chart_image(slide, chart_path, left=Inches(0.6), top=Inches(3.6), width=half)
            if spec.get("table"):
                _add_table(slide, spec["table"]["header"], spec["table"]["rows"],
                           top=Inches(1.5), left=Inches(0.9) + half, width=half)

        elif kind == "chart":
            _add_bullets(slide, bullets, top=Inches(1.3), height=Inches(2.4), size=14)
            if chart_path:
                _add_chart_image(slide, chart_path, left=Inches(2.0), top=Inches(3.9),
                                 width=Inches(9.0))
            else:
                note = slide.shapes.add_textbox(Inches(0.6), Inches(4.2), SW - Inches(1.2), Inches(0.5))
                np_ = note.text_frame.paragraphs[0]
                nr = np_.add_run()
                nr.text = "(diagram: see docs/slide_deck_outline.md mermaid source)"
                nr.font.size = Pt(11)
                nr.font.italic = True
                nr.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

        elif kind == "closing":
            _add_bullets(slide, bullets, top=Inches(1.3), height=Inches(2.6), size=15)
            examples = spec.get("examples", [])
            if examples:
                _add_bullets(slide, ["Sample questions (public demo KB):"] + examples,
                             top=Inches(4.2), height=Inches(2.2), size=13)

        else:  # plain bullets
            _add_bullets(slide, bullets, top=Inches(1.4), height=Inches(5.0), size=17)

        _add_footer(slide, idx, total)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


def _render_markdown(specs: List[Dict[str, Any]], charts: Dict[str, Path], out_path: Path) -> Path:
    """Markdown fallback deck when python-pptx is unavailable."""
    lines: List[str] = ["# KBQA Presentation (Markdown fallback)\n",
                         "_python-pptx not installed — generated a Markdown deck instead._\n"]
    for idx, spec in enumerate(specs, start=1):
        lines.append("\n---\n")
        lines.append("## Slide {} — {}\n".format(idx, spec.get("title", "")))
        if spec.get("kind") == "title" and spec.get("subtitle"):
            for line in spec["subtitle"].split("\n"):
                lines.append("> {}\n".format(line))
        for b in spec.get("bullets", []):
            lines.append("- {}\n".format(b))
        if spec.get("table"):
            header = spec["table"]["header"]
            lines.append("\n| " + " | ".join(str(h) for h in header) + " |\n")
            lines.append("|" + "|".join(["---"] * len(header)) + "|\n")
            for row in spec["table"]["rows"]:
                cells = [str(row[c]) if c < len(row) else "" for c in range(len(header))]
                lines.append("| " + " | ".join(cells) + " |\n")
        chart_path = _match_chart(charts, spec.get("chart_keys", [])) if charts else None
        if chart_path:
            lines.append("\n![{}]({})\n".format(chart_path.stem, chart_path.name))
        for ex in spec.get("examples", []):
            lines.append("- {}\n".format(ex))
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text("".join(lines), encoding="utf-8")
    return out_path


# ─────────────────────────────────────────────────────────────────────────────
# Public entrypoint
# ─────────────────────────────────────────────────────────────────────────────

def generate_slides(
    cfg: AppConfig,
    title: Optional[str] = None,
    author: Optional[str] = None,
    out_path: Optional[Any] = None,
) -> Path:
    """Build the ~12-slide KBQA PPTX deck and return its path.

    Parameters
    ----------
    cfg
        Loaded :class:`AppConfig`.
    title, author
        Override the deck title/author; default to ``cfg.project_title`` /
        ``cfg.author``.
    out_path
        Explicit output path. When ``None``, writes to
        ``artifacts_dir()/submission/submission-<stamp>/slides.pptx``.

    Returns
    -------
    pathlib.Path
        Path to the written ``.pptx`` (or ``.md`` fallback if python-pptx is
        missing). Never raises for missing optional dependencies or artifacts.
    """
    title = title or getattr(cfg, "project_title", _DECK_TITLE)
    author = author or getattr(cfg, "author", "Le Dinh Minh Quan")

    # Resolve output location.
    if out_path is not None:
        out_path = Path(out_path)
        sub_dir = out_path.parent
    else:
        sub_dir = artifacts_dir() / "submission" / "submission-{}".format(utc_stamp())
        out_path = sub_dir / "slides.pptx"
    sub_dir.mkdir(parents=True, exist_ok=True)

    # 1) Load artifacts (live eval metrics etc.) defensively.
    art = _try_load_all_artifacts(cfg)

    # 2) Render charts into the submission dir (best-effort).
    charts = _build_charts(cfg, art, sub_dir / "charts")

    # 3) Assemble slide specs from outline + live numbers.
    specs = _slide_specs(art, title, author)

    # 4) Render PPTX, or degrade to Markdown if python-pptx is missing.
    try:
        import pptx  # noqa: F401  (presence check; lazy heavy dep)
    except Exception as exc:
        md_path = out_path.with_suffix(".md")
        logger.warning("python-pptx unavailable (%s); writing Markdown fallback -> %s", exc, md_path)
        return _render_markdown(specs, charts, md_path)

    try:
        result = _render_pptx(specs, charts, out_path)
        logger.info("Slides -> %s (%d slides, %d charts)", result, len(specs), len(charts))
        return result
    except Exception as exc:  # pragma: no cover - defensive
        md_path = out_path.with_suffix(".md")
        logger.warning("PPTX render failed (%s); writing Markdown fallback -> %s", exc, md_path)
        return _render_markdown(specs, charts, md_path)
